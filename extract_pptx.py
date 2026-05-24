import collections 
import collections.abc
try:
    from pptx import Presentation
except ImportError:
    import sys
    print("python-pptx is not installed")
    sys.exit(1)

def extract_text():
    prs = Presentation("arm_based_soc_verification.pptx")
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)

if __name__ == "__main__":
    extract_text()
